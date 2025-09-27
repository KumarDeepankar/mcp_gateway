# agentic_assistant/document_parser.py
import io
import pandas as pd
import docx
from pptx import Presentation
from fastapi import HTTPException

def extract_excel_content(file_content: bytes) -> str:
    """Extract text and data from Excel files"""
    try:
        excel_file = io.BytesIO(file_content)
        xl_file = pd.ExcelFile(excel_file)
        extracted_data = []

        for sheet_name in xl_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            sheet_info = f"\n--- Sheet: {sheet_name} ---\n"
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                sheet_info += f"Numeric columns: {', '.join(numeric_cols)}\n"
                sheet_info += f"Data shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
            sheet_info += "\nSample data:\n"
            sheet_info += df.head(10).to_string(index=False)
            if len(numeric_cols) > 0:
                sheet_info += f"\n\nSummary statistics:\n"
                sheet_info += df[numeric_cols].describe().to_string()
            extracted_data.append(sheet_info)
        return "\n".join(extracted_data)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing Excel file: {e}")

def extract_word_content(file_content: bytes) -> str:
    """Extract text content from Word documents"""
    try:
        doc_file = io.BytesIO(file_content)
        doc = docx.Document(doc_file)
        extracted_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                extracted_text.append(paragraph.text.strip())
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            if table_data:
                extracted_text.append("\n--- Table Data ---")
                for row in table_data:
                    extracted_text.append(" | ".join(row))
        return "\n".join(extracted_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing Word file: {e}")

def extract_ppt_content(file_content: bytes) -> str:
    """Extract text content from PowerPoint presentations"""
    try:
        ppt_file = io.BytesIO(file_content)
        prs = Presentation(ppt_file)
        extracted_text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"\n--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                if shape.has_table:
                    table = shape.table
                    slide_content.append("\n--- Table in Slide ---")
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        slide_content.append(" | ".join(row_data))
            extracted_text.extend(slide_content)
        return "\n".join(extracted_text)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing PowerPoint file: {e}")